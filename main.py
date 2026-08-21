import os
import tkinter as tk
from tkinter import filedialog, messagebox, ttk
from PIL import Image
from PyPDF2 import PdfMerger
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
from reportlab.lib import colors
import openpyxl


class UniversalMergerApp:
    def __init__(self, root):
        self.root = root
        self.root.title("Universal File Merger")
        self.root.geometry("600x400")

        self.file_list = []

        # UI Setup
        self.label = ttk.Label(root, text="Select PDF, Excel, PPTX, or PNG files to merge:")
        self.label.pack(pady=10)

        self.listbox = tk.Listbox(root, selectmode=tk.SINGLE, width=70, height=12)
        self.listbox.pack(padx=10, pady=5)

        # Buttons Frame
        btn_frame = ttk.Frame(root)
        btn_frame.pack(pady=10)

        self.add_btn = ttk.Button(btn_frame, text="Add Files", command=self.add_files)
        self.add_btn.grid(row=0, column=0, padx=5)

        self.remove_btn = ttk.Button(btn_frame, text="Remove Selected", command=self.remove_file)
        self.remove_btn.grid(row=0, column=1, padx=5)

        self.clear_btn = ttk.Button(btn_frame, text="Clear All", command=self.clear_files)
        self.clear_btn.grid(row=0, column=2, padx=5)

        self.merge_btn = ttk.Button(root, text="Merge to PDF", command=self.merge_files)
        self.merge_btn.pack(pady=10)

    def add_files(self):
        files = filedialog.askopenfilenames(
            title="Select Files",
            filetypes=[("Supported Files", "*.pdf *.xlsx *.pptx *.png")]
        )
        for file in files:
            if file not in self.file_list:
                self.file_list.append(file)
                self.listbox.insert(tk.END, os.path.basename(file))

    def remove_file(self):
        selected_idx = self.listbox.curselection()
        if selected_idx:
            idx = selected_idx[0]
            self.listbox.delete(idx)
            del self.file_list[idx]

    def clear_files(self):
        self.file_list.clear()
        self.listbox.delete(0, tk.END)

    def excel_to_pdf(self, excel_path, temp_pdf_path):
        """Converts an Excel sheet into a temporary PDF table."""
        wb = openpyxl.load_workbook(excel_path, data_only=True)
        sheet = wb.active
        data = []

        for row in sheet.iter_rows(values_only=True):
            if any(row):  # Skip completely empty rows
                data.append([str(cell) if cell is not None else "" for cell in row])

        if not data:
            data = [["Empty Sheet"]]

        doc = SimpleDocTemplate(temp_pdf_path, pagesize=letter)
        t = Table(data)
        t.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.black),
            ('FONTSIZE', (0, 0), (-1, -1), 8),
        ]))
        doc.build([t])

    def image_to_pdf(self, img_path, temp_pdf_path):
        """Converts a PNG image to a PDF file."""
        image = Image.open(img_path)
        if image.mode != 'RGB':
            image = image.convert('RGB')
        image.save(temp_pdf_path, "PDF", resolution=100.0)

    def ppt_to_pdf_text(self, ppt_path, temp_pdf_path):
        """Extracts text content from PowerPoint slides into a PDF table."""
        prs = Presentation(ppt_path)
        data = [["Slide #", "Extracted Content"]]

        for idx, slide in enumerate(prs.slides, start=1):
            text_runs = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text_runs.append(paragraph.text)
            slide_text = "\n".join(text_runs).strip()
            data.append([f"Slide {idx}", slide_text if slide_text else "[No Text Content]"])

        doc = SimpleDocTemplate(temp_pdf_path, pagesize=letter)
        t = Table(data, colWidths=[80, 400])
        t.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.navy),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
            ('VALIGN', (0, 0), (-1, -1), 'TOP'),
            ('FONTSIZE', (0, 0), (-1, -1), 9),
        ]))
        doc.build([t])

    def merge_files(self):
        if not self.file_list:
            messagebox.showwarning("Warning", "Please add at least one file to merge.")
            return

        save_path = filedialog.asksaveasfilename(
            defaultextension=".pdf",
            filetypes=[("PDF Documents", "*.pdf")]
        )
        if not save_path:
            return

        merger = PdfMerger()
        temp_files = []

        try:
            for i, filepath in enumerate(self.file_list):
                ext = os.path.splitext(filepath)[1].lower()

                if ext == ".pdf":
                    merger.append(filepath)
                elif ext == ".png":
                    temp_pdf = f"temp_img_{i}.pdf"
                    self.image_to_pdf(filepath, temp_pdf)
                    temp_files.append(temp_pdf)
                    merger.append(temp_pdf)
                elif ext == ".xlsx":
                    temp_pdf = f"temp_excel_{i}.pdf"
                    self.excel_to_pdf(filepath, temp_pdf)
                    temp_files.append(temp_pdf)
                    merger.append(temp_pdf)
                elif ext == ".pptx":
                    temp_pdf = f"temp_ppt_{i}.pdf"
                    self.ppt_to_pdf_text(filepath, temp_pdf)
                    temp_files.append(temp_pdf)
                    merger.append(temp_pdf)

            merger.write(save_path)
            merger.close()
            messagebox.showinfo("Success", f"Files successfully merged to:\n{save_path}")

        except Exception as e:
            messagebox.showerror("Error", f"An error occurred during merging:\n{str(e)}")

        finally:
            # Clean up temporary PDF files
            for temp_file in temp_files:
                if os.path.exists(temp_file):
                    os.remove(temp_file)


if __name__ == "__main__":
    root = tk.Tk()
    app = UniversalMergerApp(root)
    root.mainloop()